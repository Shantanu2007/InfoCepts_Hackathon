import streamlit as st
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt
from pathlib import Path
import tempfile
import zipfile

# =====================================================
# PAGE CONFIG
# =====================================================

st.set_page_config(
    page_title="Revenue Reporting Agent",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded"
)

# =====================================================
# THEME TOGGLE
# =====================================================

with st.sidebar:
    dark_mode = st.toggle(
        "🌙 Dark Mode",
        value=True
    )

# =====================================================
# THEME COLORS
# =====================================================

if dark_mode:

    BG_COLOR = "#07111F"
    SIDEBAR_COLOR = "#0B1324"
    CARD_COLOR = "#111827"

    TEXT_COLOR = "#FFFFFF"
    SUBTEXT_COLOR = "#CBD5E1"
    SIDEBAR_TEXT = "#FFFFFF"

    BORDER_COLOR = "#334155"

    HEADER_GRADIENT = """
    linear-gradient(
        90deg,
        #38BDF8,
        #6366F1,
        #A855F7
    )
    """

else:

    BG_COLOR = "#F8FAFC"
    SIDEBAR_COLOR = "#E2E8F0"
    CARD_COLOR = "#FFFFFF"

    TEXT_COLOR = "#111827"
    SUBTEXT_COLOR = "#475569"
    SIDEBAR_TEXT = "#111827"

    BORDER_COLOR = "#CBD5E1"

    HEADER_GRADIENT = """
    linear-gradient(
        90deg,
        #0284C7,
        #4F46E5,
        #7C3AED
    )
    """

# =====================================================
# CSS
# =====================================================

st.markdown(
f"""
<style>

/* Keep sidebar expand/collapse working */
header {{
    background: transparent !important;
}}

[data-testid="stHeader"] {{
    background: transparent !important;
}}

#MainMenu {{
    visibility:hidden;
}}

footer {{
    visibility:hidden;
}}

.block-container {{
    padding-top: 1rem !important;
}}

/* App */

.stApp {{
    background:{BG_COLOR};
}}

/* Sidebar */

section[data-testid="stSidebar"] {{
    background:{SIDEBAR_COLOR};
    width:215px !important;
}}

section[data-testid="stSidebar"] * {{
    color:{SIDEBAR_TEXT} !important;
}}

/* Header */

.main-title {{

    text-align:center;
    font-size:58px;
    font-weight:800;

    background:{HEADER_GRADIENT};

    -webkit-background-clip:text;
    -webkit-text-fill-color:transparent;

    margin-bottom:5px;
}}

.sub-title {{

    text-align:center;
    font-size:20px;

    color:{SUBTEXT_COLOR};

    margin-bottom:10px;
}}

/* Upload Area */

[data-testid="stFileUploader"] {{

    background:{CARD_COLOR};

    border:3px dashed #38BDF8;

    border-radius:24px;

    padding:25px;

    box-shadow:
    0px 0px 20px rgba(
    56,189,248,.15
    );
}}

[data-testid="stFileUploaderDropzone"] {{

    background:{CARD_COLOR} !important;
}}

[data-testid="stFileUploaderDropzone"] * {{

    color:{TEXT_COLOR} !important;
}}

[data-testid="stFileUploader"] * {{

    color:{TEXT_COLOR} !important;
}}

[data-testid="stBaseButton-secondary"] {{

    background:{CARD_COLOR} !important;

    color:{TEXT_COLOR} !important;

    border:1px solid {BORDER_COLOR} !important;
}}

[data-testid="stBaseButton-secondary"] span {{

    color:{TEXT_COLOR} !important;
}}

/* KPI Cards */

[data-testid="metric-container"] {{

    background:{CARD_COLOR};

    border-radius:20px;

    padding:20px;

    border:1px solid rgba(
    99,102,241,.35
    );

    box-shadow:
    0px 0px 15px rgba(
    99,102,241,.15
    );

    transition:.3s ease;
}}

[data-testid="metric-container"]:hover {{

    transform:translateY(-6px);

    box-shadow:
    0px 10px 25px rgba(
    99,102,241,.35
    );
}}

/* Buttons */

.stButton button {{

    width:100%;
    height:60px;

    border:none;

    border-radius:15px;

    font-size:18px;
    font-weight:bold;

    background:
    linear-gradient(
    90deg,
    #2563EB,
    #7C3AED
    );

    color:white;
}}

.stButton button:hover {{

    transform:scale(1.02);
}}

/* Charts */

.chart-box {{

    background:{CARD_COLOR};

    border:1px solid {BORDER_COLOR};

    border-radius:20px;

    padding:20px;
}}

h1,h2,h3,h4,h5,h6,
p,span,label,div {{

    color:{TEXT_COLOR};
}}

</style>
""",
unsafe_allow_html=True
)

# =====================================================
# SIDEBAR
# =====================================================

with st.sidebar:

    st.markdown(
        "<h1>🤖 Revenue AI Agent</h1>",
        unsafe_allow_html=True
    )

    st.success("🟢 System Online")

    st.markdown("---")

    st.markdown("### Capabilities")

    st.write("✅ Agency Segregation")
    st.write("✅ Revenue Analytics")
    st.write("✅ KPI Dashboard")
    st.write("✅ Excel Generation")
    st.write("🔜 PPT Generation")
    st.write("🔜 AI Insights")

# =====================================================
# HEADER
# =====================================================

st.markdown(
"""
<div class='main-title'>
Revenue Reporting Agent
</div>
""",
unsafe_allow_html=True
)

st.markdown(
"""
<div class='sub-title'>
AI Powered Revenue Automation Platform
</div>
""",
unsafe_allow_html=True
)

# =====================================================
# FILE UPLOAD
# =====================================================

uploaded_file = st.file_uploader(
    "📁 Upload Revenue Excel File",
    type=["xlsx"]
)

# =====================================================
# PROCESSING
# =====================================================

if uploaded_file is not None:

    try:

        df = pd.read_excel(
            uploaded_file,
            engine="openpyxl"
        )

        c1, c2, c3 = st.columns(3)

        with c1:
            st.metric(
                "📄 Records",
                f"{len(df):,}"
            )

        with c2:
            st.metric(
                "🏢 Agencies",
                df["Grandparent_Agency"].nunique()
            )

        with c3:
            st.metric(
                "📢 Advertisers",
                df["Advertiser"].nunique()
            )

        st.markdown("---")

        st.subheader("📋 Data Preview")

        st.dataframe(
            df.head(10),
            use_container_width=True
        )

        if st.button("🚀 Generate Reports"):

            output_dir = Path(
                tempfile.mkdtemp()
                )

            summary = []

            grouped = list(
                df.groupby(
                    "Grandparent_Agency"
                )
            )

            progress = st.progress(0)

            for idx, (agency, agency_df) in enumerate(grouped):

                safe_name = (
                    str(agency)
                    .replace("/", "-")
                    .replace("\\", "-")
                )

                agency_df.to_excel(
                    output_dir / f"{safe_name}.xlsx",
                    index=False,
                    engine="openpyxl"
                )

                ppt = Presentation()
                
                total_revenue = agency_df["Revenue_USD"].sum()
                
                avg_pacing = agency_df["Pacing_Pct"].mean()
                
                campaign_count = (
                        agency_df["Campaign_Name"]
                                .nunique()
                                )
                
                advertiser_count = (
                     agency_df["Advertiser"]
                     .nunique()
                     )
                
                slide = ppt.slides.add_slide(
                     ppt.slide_layouts[0]
                     )
                
                slide.shapes.title.text = (
                        f"{agency} Revenue Report"
                        )
                
                slide.placeholders[1].text = (
                        "Executive Performance Summary"
                        )
                
                textbox = slide.shapes.add_textbox(
                     Inches(1),
                    Inches(2),
                     Inches(6),
                    Inches(2)
                    )
                
                textbox.text_frame.text = f"""
                Total Revenue: ${total_revenue:,.0f}

                Average Pacing: {avg_pacing:.1f}%

                Campaigns: {campaign_count}

                Advertisers: {advertiser_count}
"""
                
                slide = ppt.slides.add_slide(
                     ppt.slide_layouts[1]
                     )
                
                slide.shapes.title.text = (
                     "Revenue Summary"
                     )
                
                top_advertiser = (
                     agency_df.groupby(
                          "Advertiser"
                          )["Revenue_USD"]
                          .sum()
                          .idxmax()
                          )
                
                
                content = slide.placeholders[1]
                
                content.text = f"""
                Total Revenue:
                ${total_revenue:,.0f}

                Top Advertiser:
                {top_advertiser}

                Active Campaigns:
                {campaign_count}
                
                Average Pacing:
                {avg_pacing:.1f}%
"""

                advertiser_revenue = (
                     agency_df
                     .groupby("Advertiser")
                     ["Revenue_USD"]
                     .sum()
                     .sort_values(
                         ascending=False
                           )
                    .head(5)
                    )
                
                plt.figure(
                     figsize=(6,4)
                     )

                advertiser_revenue.plot(
                    kind="bar"
                )

                slide = ppt.slides.add_slide(
                    ppt.slide_layouts[5]
                )
                
                slide.shapes.title.text = (
                    "Revenue by Advertiser"
                )
                
                advertiser_revenue = (
                    agency_df.groupby("Advertiser")["Revenue_USD"]
                    .sum()
                    .sort_values(ascending=False)
                    .head(5)
                )
                
                advertiser_chart_path = (
                    output_dir /
                    f"{safe_name}_advertiser.png"
                )
                
                plt.figure(figsize=(6,4))
                
                advertiser_revenue.plot(
                    kind="bar",
                    color="steelblue"
                )
                
                plt.title("Top Advertisers")
                
                plt.tight_layout()
                
                plt.savefig(advertiser_chart_path)
                
                plt.close()

                campaign_revenue = (
                    agency_df
                    .groupby("Campaign_Name")
                    ["Revenue_USD"]
                    .sum()
                    .sort_values(
                        ascending=False
                    )
                    .head(5)
                )
                
                plt.figure(
                    figsize=(6,4)
                )
                
                campaign_revenue.plot(
                    kind="barh"
                )
                
                plt.title(
                    "Top Campaigns"
                )
                
                plt.tight_layout()
                
                campaign_chart = (
                    output_dir /
                    f"{safe_name}_campaign.png"
                )
                
                plt.savefig(
                    campaign_chart
                )
                
                plt.close()
                
                slide = ppt.slides.add_slide(
                    ppt.slide_layouts[5]
                )
                
                slide.shapes.title.text = (
                    "Campaign Performance"
                )
                
                slide.shapes.add_picture(
                    str(campaign_chart),
                    Inches(1),
                    Inches(1.3),
                    width=Inches(7)
                )
                
                slide = ppt.slides.add_slide(
                    ppt.slide_layouts[1]
                )
                
                slide.shapes.title.text = (
                    "Key Takeaways"
                )
                
                content = slide.placeholders[1]
                
                content.text = f"""
                • Total Revenue generated:
                ${total_revenue:,.0f}
                
                • Top Advertiser:
                {top_advertiser}
                
                • Average Pacing:
                {avg_pacing:.1f}%
                
                • Active Campaigns:
                {campaign_count}
                
                • Advertisers:
                {advertiser_count}
                """

                plt.title(
                    "Top Advertisers"
                )
                
                plt.tight_layout()
                
                chart_path = (
                    output_dir /
                    f"{safe_name}_advertiser.png"
                )
                
                plt.savefig(chart_path)
                
                plt.close()
                
                ppt_path = (
                     output_dir /
                     f"{safe_name}.pptx"
                     )
                ppt.save(str(ppt_path))

                
                content = slide.placeholders[1]
                
                content.text = f"""
                Total Revenue:
                ${agency_df['Revenue_USD'].sum():,.0f}
                
                Average Pacing:
                {agency_df['Pacing_Pct'].mean():.2f}%
                Campaign Count:
                {agency_df['Campaign_Name'].nunique()}

                
                
                Advertiser Count:
                {agency_df['Advertiser'].nunique()}
"""

                summary.append(
                    {
                        "Agency": agency,
                        "Total Revenue": round(
                            agency_df["Revenue_USD"].sum(),
                            2
                        ),
                        "Average Pacing": round(
                            agency_df["Pacing_Pct"].mean(),
                            2
                        ),
                        "Campaign Count":
                        agency_df["Campaign_Name"].nunique(),
                        "Advertiser Count":
                        agency_df["Advertiser"].nunique()
                    }
                )

                progress.progress(
                    (idx + 1) / len(grouped)
                )

            summary_df = pd.DataFrame(
                summary
            ).sort_values(
                "Total Revenue",
                ascending=False
            )

            summary_df.to_excel(
                output_dir / "Agency_Summary.xlsx",
                index=False,
                engine="openpyxl"
            )

            zip_file_path = (
                output_dir /
                "RevenueReports.zip"
                )
            
            with zipfile.ZipFile(
                zip_file_path,
                "w",
                zipfile.ZIP_DEFLATED
                ) as zipf:
                
                for file in output_dir.iterdir():
                    if file.suffix in [
                        ".xlsx",
                        ".pptx"
                        ]:
                        
                        zipf.write(
                            file,
                            arcname=file.name
                            )

            total_revenue = (
                summary_df["Total Revenue"].sum()
            )

            top_agency = (
                summary_df.iloc[0]["Agency"]
            )

            a1, a2 = st.columns(2)

            with a1:
                st.metric(
                    "💰 Portfolio Revenue",
                    f"${total_revenue:,.0f}"
                )

            with a2:
                st.metric(
                    "🏆 Top Agency",
                    top_agency
                )

            chart1, chart2 = st.columns(2)

            with chart1:
                st.subheader("💰 Revenue By Agency")

                st.bar_chart(
                    summary_df.set_index(
                        "Agency"
                    )["Total Revenue"]
                )

            with chart2:
                st.subheader("📈 Average Pacing")

                st.line_chart(
                    summary_df.set_index(
                        "Agency"
                    )["Average Pacing"]
                )

            st.subheader("📊 Agency Summary")

            st.dataframe(
                summary_df,
                use_container_width=True
            )

            st.success(
                 "🎉 Reports Generated Successfully"
                 )
            
            with open(
                zip_file_path,
                "rb"
                ) as file:
                
                st.download_button(
                    label="📥 Download All Reports",
                    data=file,
                    file_name="RevenueReports.zip",
                    mime="application/zip",
                    use_container_width=True
                    )

    except Exception as e:

        st.error(
            f"Error: {str(e)}"
        )
